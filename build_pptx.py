"""
Build ECG Interpreter PPTX — Dr. Femi Lawal · RAIN · March 27 2026
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as u
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0f, 0x4c, 0x81)
TEAL   = RGBColor(0x1a, 0x9e, 0x8f)
CORAL  = RGBColor(0xe6, 0x39, 0x46)
AMBER  = RGBColor(0xf4, 0xa2, 0x61)
WHITE  = RGBColor(0xff, 0xff, 0xff)
LTGRAY = RGBColor(0xf8, 0xfa, 0xfc)
MUTED  = RGBColor(0x64, 0x74, 0x8b)
DARK   = RGBColor(0x1e, 0x29, 0x3b)
BORDER = RGBColor(0xe2, 0xe8, 0xf0)

prs = Presentation()
prs.slide_width  = Inches(13.33)   # 16:9 widescreen
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helper functions ──────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def grad_bg(slide, c1=NAVY, c2=TEAL):
    """Approximated gradient with solid navy (PPT gradients are complex)"""
    fill = slide.background.fill
    fill.gradient()
    fill.gradient_angle = 315
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[0].position  = 0.0
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_stops[1].position  = 1.0

def txb(slide, text, l, t, w, h,
        size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    tf = slide.shapes.add_textbox(l, t, w, h)
    tf.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf

def rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shp = slide.shapes.add_shape(
        pptx.util.MSO_SHAPE_TYPE if False else 1,   # freeform → use add_shape
        l, t, w, h)
    # use add_shape MSO_CONNECTOR_TYPE.STRAIGHT isn't right; use add_shape with rectange type
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shp.fill.solid() if fill else shp.fill.background()
    if fill:
        shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        if line_w:
            shp.line.width = line_w
    else:
        shp.line.fill.background()
    return shp

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_pt=0.75):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Pt as Pt2
    shp = slide.shapes.add_shape(1, l, t, w, h)  # 1 = rectangle
    if fill_rgb:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_rgb
    else:
        shp.fill.background()
    if line_rgb:
        shp.line.color.rgb = line_rgb
        shp.line.width = Pt2(line_pt)
    else:
        shp.line.fill.background()
    return shp

def card(slide, l, t, w, h, fill=WHITE, accent=NAVY, title="", bullets=None, title_size=14, bullet_size=11):
    """White card with colored top accent bar"""
    add_rect(slide, l, t, w, h, fill_rgb=fill, line_rgb=BORDER, line_pt=0.5)
    # accent top bar
    add_rect(slide, l, t, w, Inches(0.07), fill_rgb=accent)
    if title:
        txb(slide, title, l+Inches(0.1), t+Inches(0.1), w-Inches(0.2), Inches(0.35),
            size=title_size, bold=True, color=accent)
    if bullets:
        yt = t + Inches(0.48)
        for b in bullets:
            tf = slide.shapes.add_textbox(l+Inches(0.15), yt, w-Inches(0.3), Inches(0.3))
            tf.word_wrap = True
            p = tf.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "• " + b
            run.font.size = Pt(bullet_size)
            run.font.color.rgb = DARK
            yt += Inches(0.25)

def metric_card(slide, l, t, w, h, value, label, accent=NAVY):
    add_rect(slide, l, t, w, h, fill_rgb=WHITE, line_rgb=BORDER, line_pt=0.5)
    add_rect(slide, l, t, w, Inches(0.07), fill_rgb=accent)
    txb(slide, value, l, t+Inches(0.1), w, Inches(0.45),
        size=24, bold=True, color=accent, align=PP_ALIGN.CENTER)
    txb(slide, label, l, t+Inches(0.55), w, Inches(0.35),
        size=9, color=MUTED, align=PP_ALIGN.CENTER)

def tag_box(slide, text, l, t, color=NAVY):
    w = Inches(1.8); h = Inches(0.28)
    add_rect(slide, l, t, w, h, fill_rgb=color)
    txb(slide, text.upper(), l, t, w, h, size=8, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)

def section_title(slide, text, l=Inches(0.4), t=Inches(0.65), color=NAVY):
    tf = slide.shapes.add_textbox(l, t, W - l - Inches(0.4), Inches(0.55))
    p = tf.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(24)
    run.font.bold  = True
    run.font.color.rgb = color
    # underline via bottom border rect
    add_rect(slide, l, t + Inches(0.48), Inches(3.5), Inches(0.05), fill_rgb=TEAL)

def bullet_list(slide, items, l, t, w, h, size=11, color=DARK, indent=True):
    tf = slide.shapes.add_textbox(l, t, w, h)
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.text_frame.paragraphs[0]
        else:
            p = tf.text_frame.add_paragraph()
        run = p.add_run()
        run.text = ("• " if indent else "") + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_before = Pt(2)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
grad_bg(sl, NAVY, TEAL)

# Institution label
add_rect(sl, Inches(4.0), Inches(0.35), Inches(5.33), Inches(0.38),
         fill_rgb=RGBColor(0xff,0xff,0xff) if False else None)
inst_box = sl.shapes.add_shape(1, Inches(4.0), Inches(0.3), Inches(5.33), Inches(0.42))
inst_box.fill.solid(); inst_box.fill.fore_color.rgb = RGBColor(0x1a,0x6a,0xaa)
inst_box.line.fill.background()
txb(sl, "AI  ECG  INTERPRETER", Inches(4.0), Inches(0.3), Inches(5.33), Inches(0.42),
    size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Title
txb(sl, "ECG Interpreter", Inches(0.5), Inches(0.85), W-Inches(1.0), Inches(1.0),
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
txb(sl, "AI-Powered Cardiac Diagnostics for Resource-Limited Settings",
    Inches(0.5), Inches(1.85), W-Inches(1.0), Inches(0.5),
    size=18, color=RGBColor(0xcc,0xe8,0xff), align=PP_ALIGN.CENTER)

# Metric cards
metrics = [("0.954","PTB-XL AUROC"), ("0.833","SHD AUROC"),
           ("21,799","Training ECGs"), ("24","Conditions Detected")]
mw = Inches(2.6); mh = Inches(1.1); mg = Inches(0.18)
ml = (W - 4*mw - 3*mg) / 2
for i,(v,lb) in enumerate(metrics):
    accent = [TEAL, CORAL, AMBER, RGBColor(0x6d,0x68,0x75)][i]
    bx = sl.shapes.add_shape(1, ml+i*(mw+mg), Inches(2.55), mw, mh)
    bx.fill.solid(); bx.fill.fore_color.rgb = RGBColor(0x1a,0x6a,0xaa)
    bx.line.color.rgb = RGBColor(0xff,0xff,0xff); bx.line.width = Pt(0.5)
    txb(sl, v,  ml+i*(mw+mg), Inches(2.62), mw, Inches(0.55),
        size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, lb, ml+i*(mw+mg), Inches(3.18), mw, Inches(0.38),
        size=10, color=RGBColor(0xaa,0xcc,0xff), align=PP_ALIGN.CENTER)

# ECG wave decoration (simple line)
add_rect(sl, Inches(0.4), Inches(3.88), W-Inches(0.8), Inches(0.03), fill_rgb=RGBColor(0xff,0xff,0xff))

# Author / date at bottom
txb(sl, "Dr. Femi Lawal  ·  Robotic and Artificial Intelligence Nigeria (RAIN)  ·  March 27, 2026",
    Inches(0.4), H-Inches(0.55), W-Inches(0.8), Inches(0.42),
    size=11, color=RGBColor(0xaa,0xcc,0xff), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — NIGERIA PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl, WHITE)
tag_box(sl, "The Challenge", Inches(0.4), Inches(0.2), CORAL)
section_title(sl, "Nigeria's Cardiac Care Gap")

# Metric cards (2×2)
pairs = [("1:10K","Cardiologist-to-patient ratio [3]",CORAL),
         ("<300", "Echo machines nationwide [2]",    AMBER),
         ("60%",  "CVD deaths in LMICs — preventable [1]", CORAL),
         ("216M", "Nigeria's population [1]", TEAL)]
cw = Inches(2.8); ch = Inches(1.35); cg = Inches(0.15)
for i,(v,lb,ac) in enumerate(pairs):
    row, col = divmod(i, 2)
    cx = Inches(0.4) + col*(cw+cg)
    cy = Inches(1.35) + row*(ch+cg)
    metric_card(sl, cx, cy, cw, ch, v, lb, ac)

# Right column boxes
rx = Inches(6.3); rw = W - rx - Inches(0.35)

# Coral box
add_rect(sl, rx, Inches(1.35), Inches(0.08), Inches(1.45), fill_rgb=CORAL)
add_rect(sl, rx+Inches(0.08), Inches(1.35), rw-Inches(0.08), Inches(1.45),
         fill_rgb=RGBColor(0xfe,0xf2,0xf2), line_rgb=BORDER)
txb(sl, "The Diagnostic Gap", rx+Inches(0.2), Inches(1.42), rw-Inches(0.3), Inches(0.32),
    size=12, bold=True, color=CORAL)
bullet_list(sl,
    ["Most patients diagnosed late or not at all",
     "Specialists concentrated in Lagos & Abuja",
     "Referral chains are long and costly [2]"],
    rx+Inches(0.2), Inches(1.78), rw-Inches(0.3), Inches(0.85), size=10)

# Teal box
add_rect(sl, rx, Inches(2.9), Inches(0.08), Inches(1.65), fill_rgb=TEAL)
add_rect(sl, rx+Inches(0.08), Inches(2.9), rw-Inches(0.08), Inches(1.65),
         fill_rgb=RGBColor(0xf0,0xfa,0xf8), line_rgb=BORDER)
txb(sl, "The Opportunity", rx+Inches(0.2), Inches(2.97), rw-Inches(0.3), Inches(0.32),
    size=12, bold=True, color=TEAL)
bullet_list(sl,
    ["ECG machines in rural PHCs (₦480K–₦1.6M) [8]",
     "ECG costs ₦4,800–₦12,800 vs ₦320K–₦640K for echo [8]",
     "AI bridges the skills gap without new hardware"],
    rx+Inches(0.2), Inches(3.33), rw-Inches(0.3), Inches(1.1), size=10)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl, WHITE)
tag_box(sl, "The Solution", Inches(0.4), Inches(0.2), TEAL)
section_title(sl, "ECG Interpreter — What It Does")

boxes = [
    (TEAL,  "Paper ECG Friendly",
     ["Accepts smartphone photos of printed ECGs",
      "AI digitiser converts image → waveform",
      "Works with any 12-lead machine"]),
    (NAVY,  "Fully Offline",
     ["Runs on a laptop or tablet",
      "No internet or cloud subscription needed",
      "~121 MB total model weight"]),
    (CORAL, "Clinical Decision Support",
     ["Flags HIGH / MOD / LOW risk in <30 s",
      "Evidence-weighted differential diagnoses",
      "Augments — does not replace — the clinician"]),
]
bw = Inches(3.8); bg_h = Inches(2.4); bg_g = Inches(0.25)
bl = (W - 3*bw - 2*bg_g)/2
for i,(ac,title,buls) in enumerate(boxes):
    bx = bl + i*(bw+bg_g)
    card(sl, bx, Inches(1.42), bw, bg_h, WHITE, ac, title, buls, title_size=14, bullet_size=11)

# Amber highlight box
add_rect(sl, Inches(0.4), Inches(4.1), W-Inches(0.8), Inches(0.9),
         fill_rgb=RGBColor(0xff,0xfb,0xf0), line_rgb=AMBER)
txb(sl, "Real scenario: PHC nurse photographs ECG → AI flags Reduced EF + Pulmonary HTN (HIGH risk) → "
        "Urgent referral issued. Cost: ₦4,800 ECG + a smartphone photo. [8]",
    Inches(0.6), Inches(4.18), W-Inches(1.2), Inches(0.75),
    size=11, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl, WHITE)
tag_box(sl, "Architecture", Inches(0.4), Inches(0.2), NAVY)
section_title(sl, "End-to-End Pipeline")

# Pipeline boxes
steps = [("1. Input","JPEG/PNG/PDF"),("2. Digitiser","Image→Waveform"),
         ("3. Features","NeuroKit2 [7]"),("4. Dual AI","Two Neural Nets"),
         ("5. Report","Risk+Diagnoses")]
sw = Inches(2.0); sh = Inches(0.65); sg = Inches(0.08)
sl_x = (W - 5*sw - 4*(sg+Inches(0.25)))/2
for i,(t1,t2) in enumerate(steps):
    px = sl_x + i*(sw + sg + Inches(0.25))
    add_rect(sl, px, Inches(1.42), sw, sh, fill_rgb=LTGRAY, line_rgb=BORDER, line_pt=1)
    txb(sl, t1, px, Inches(1.46), sw, Inches(0.3), size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txb(sl, t2, px, Inches(1.72), sw, Inches(0.28), size=9, color=MUTED, align=PP_ALIGN.CENTER)
    if i < 4:
        txb(sl, "→", px+sw+sg, Inches(1.55), Inches(0.25), Inches(0.3),
            size=16, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Two model boxes
mw2 = (W - Inches(1.1))/2
# Model 1
card(sl, Inches(0.4), Inches(2.25), mw2, Inches(2.75), WHITE, NAVY,
     "Model 1 — PTB-XL ECG Classifier [4]", title_size=12)
layers1 = ["Input: 12-lead ECG · 12×2,500 samples @ 250 Hz",
           "1D ResNet Backbone (ECGEncoder1D) — 4 stages, ~5M params",
           "256-dim ECG Embedding (Global Avg Pool)",
           "→ 12 ECG Diagnoses (sigmoid output)",
           "PTB-XL · 21,799 records · Val AUROC 0.954"]
bullet_list(sl, layers1, Inches(0.55), Inches(2.75), mw2-Inches(0.3), Inches(2.1), 10)

# Model 2
card(sl, Inches(0.55)+mw2, Inches(2.25), mw2, Inches(2.75), WHITE, CORAL,
     "Model 2 — EchoNext Structural Detector [6]", title_size=12)
layers2 = ["ECG Branch → 256-dim  |  Tabular (7 features) → 128-dim",
           "Fusion MLP: 384 → 256 → 64 (concatenated)",
           "Dropout 0.2 for regularisation",
           "→ 12 Structural Conditions (sigmoid output)",
           "Columbia CUIMC · Echo-confirmed labels · SHD AUROC 0.833"]
bullet_list(sl, layers2, Inches(0.7)+mw2, Inches(2.75), mw2-Inches(0.3), Inches(2.1), 10)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DATASETS
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl, WHITE)
tag_box(sl, "Data", Inches(0.4), Inches(0.2), TEAL)
section_title(sl, "Training Datasets")

dw = (W - Inches(1.1))/2

# PTB-XL card
card(sl, Inches(0.4), Inches(1.35), dw, Inches(3.5), WHITE, NAVY,
     "PTB-XL v1.0.3  [4]", title_size=13)
rows1 = [("Source","PhysioNet / Univ. Basel"),("Total records","21,799"),
         ("Train / Val / Test","17,439 / 2,180 / 2,180"),("Sampling rate","100 Hz → resampled 250 Hz"),
         ("Labels","SCP codes — dual cardiologist verified"),("Size","601 MB")]
yt = Inches(1.82)
for lbl,val in rows1:
    add_rect(sl, Inches(0.45), yt, dw-Inches(0.1), Inches(0.3),
             fill_rgb=LTGRAY if rows1.index((lbl,val))%2==0 else WHITE, line_rgb=BORDER, line_pt=0.25)
    txb(sl, lbl, Inches(0.55), yt, Inches(2.0), Inches(0.28), size=10, color=MUTED)
    txb(sl, val, Inches(2.6), yt, dw-Inches(2.3), Inches(0.28), size=10, bold=True, color=DARK)
    yt += Inches(0.31)
# teal note
add_rect(sl, Inches(0.45), yt+Inches(0.05), dw-Inches(0.1), Inches(0.45),
         fill_rgb=RGBColor(0xf0,0xfa,0xf8), line_rgb=TEAL)
txb(sl, "12 binary classes: MI territories · BBB · AF · LVH · ST-T changes · Normal",
    Inches(0.6), yt+Inches(0.08), dw-Inches(0.35), Inches(0.4), size=10, color=DARK)

# EchoNext card
card(sl, Inches(0.55)+dw, Inches(1.35), dw, Inches(3.5), WHITE, CORAL,
     "EchoNext (CUIMC)  [6]", title_size=13)
rows2 = [("Source","Columbia University Medical Centre"),("Ground truth","Echocardiography findings"),
         ("ECG shape","(N, 1, 2500, 12) @ 250 Hz"),("Tabular features","Sex · HR · PR · QRS · QTc · Age"),
         ("Labels","12 binary structural conditions"),("Best checkpoint","61 MB")]
yt2 = Inches(1.82); rx2 = Inches(0.7)+dw
for lbl,val in rows2:
    add_rect(sl, rx2, yt2, dw-Inches(0.1), Inches(0.3),
             fill_rgb=LTGRAY if rows2.index((lbl,val))%2==0 else WHITE, line_rgb=BORDER, line_pt=0.25)
    txb(sl, lbl, rx2+Inches(0.1), yt2, Inches(2.0), Inches(0.28), size=10, color=MUTED)
    txb(sl, val, rx2+Inches(2.15), yt2, dw-Inches(2.3), Inches(0.28), size=10, bold=True, color=DARK)
    yt2 += Inches(0.31)
add_rect(sl, rx2, yt2+Inches(0.05), dw-Inches(0.1), Inches(0.45),
         fill_rgb=RGBColor(0xff,0xfb,0xf0), line_rgb=AMBER)
txb(sl, "Multi-modal: ECG waveform + 7 clinical features → +3–5% AUROC gain",
    rx2+Inches(0.15), yt2+Inches(0.08), dw-Inches(0.35), Inches(0.4), size=10, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TRAINING CONFIGURATION
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl, WHITE)
tag_box(sl, "Training", Inches(0.4), Inches(0.2), NAVY)
section_title(sl, "Training Configuration")

# Left: shared params
card(sl, Inches(0.4), Inches(1.35), Inches(5.8), Inches(3.45), WHITE, NAVY,
     "Shared Hyperparameters", title_size=13)
shared = [("Optimizer","AdamW"),("Learning rate","3×10⁻⁴ · cosine annealing"),
          ("Weight decay","1×10⁻⁴"),("Loss function","BCEWithLogitsLoss (class-weighted)"),
          ("Hardware","Apple M-series · MPS GPU")]
yt = Inches(1.82)
for lbl,val in shared:
    add_rect(sl, Inches(0.45), yt, Inches(5.7), Inches(0.3),
             fill_rgb=LTGRAY if shared.index((lbl,val))%2==0 else WHITE, line_rgb=BORDER, line_pt=0.25)
    txb(sl, lbl, Inches(0.55), yt, Inches(2.2), Inches(0.28), size=10, color=MUTED)
    txb(sl, val, Inches(2.8), yt, Inches(3.1), Inches(0.28), size=10, bold=True, color=DARK)
    yt += Inches(0.31)

# Time cards
for i,(v,lb,ac) in enumerate([("~8 h","EchoNext (15 epochs)",TEAL),("~3 h","PTB-XL (17 epochs)",NAVY)]):
    metric_card(sl, Inches(0.45)+i*Inches(2.92), Inches(3.25), Inches(2.75), Inches(1.1), v, lb, ac)

# Right: model-specific + augmentation
rx = Inches(6.45); rw = W - rx - Inches(0.35)
card(sl, rx, Inches(1.35), rw, Inches(1.85), WHITE, TEAL, "Model-Specific Settings", title_size=12)
tbl_data = [["","PTB-XL","EchoNext"],["Batch size","64","32"],["Max epochs","50","30"],
            ["Early stopping","patience 10","patience 7"],["Best epoch","7","7"]]
yt = Inches(1.82)
for row in tbl_data:
    bg_c = NAVY if row[0]=="" else (LTGRAY if tbl_data.index(row)%2==0 else WHITE)
    add_rect(sl, rx+Inches(0.05), yt, rw-Inches(0.1), Inches(0.28),
             fill_rgb=bg_c, line_rgb=BORDER, line_pt=0.25)
    ww = (rw-Inches(0.1))/3
    for j,cell in enumerate(row):
        txb(sl, cell, rx+Inches(0.05)+j*ww, yt, ww, Inches(0.26),
            size=9, bold=(row[0]==""), color=WHITE if row[0]=="" else DARK,
            align=PP_ALIGN.CENTER)
    yt += Inches(0.29)

card(sl, rx, Inches(3.3), rw, Inches(1.7), WHITE, CORAL, "Data Augmentation (EchoNext)", title_size=12)
bullet_list(sl,
    ["Gaussian noise — σ ≈ 1% of signal amplitude",
     "Amplitude scaling — ×0.9 to ×1.1",
     "Baseline wander — low-frequency drift",
     "Class weighting — up-weights rare conditions"],
    rx+Inches(0.15), Inches(3.78), rw-Inches(0.3), Inches(1.1), 10)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TRAINING CURVES (static table since no chart rendering)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl, WHITE)
tag_box(sl, "Results", Inches(0.4), Inches(0.2), TEAL)
section_title(sl, "Training Curves — Key Epoch Data")

# PTB-XL table
card(sl, Inches(0.4), Inches(1.35), Inches(5.8), Inches(3.9), WHITE, NAVY, "PTB-XL Classifier", title_size=13)
ptb_rows = [["Epoch","Train Loss","Val AUROC","Note"],
            ["1","0.631","0.914","Baseline"],["3","0.417","0.940","Improving"],
            ["5","0.371","0.939","Stable"],["7 ★","0.342","0.954","BEST"],
            ["10","0.309","0.952","Plateau"],["17","0.250","0.948","Final"]]
yt = Inches(1.82)
for row in ptb_rows:
    hdr = row[0]=="Epoch"
    bc  = NAVY if hdr else (RGBColor(0xd1,0xfa,0xe5) if "★" in row[0] else (LTGRAY if ptb_rows.index(row)%2==0 else WHITE))
    add_rect(sl, Inches(0.45), yt, Inches(5.7), Inches(0.3),
             fill_rgb=bc, line_rgb=BORDER, line_pt=0.25)
    ww = Inches(5.7)/4
    for j,cell in enumerate(row):
        txb(sl, cell, Inches(0.45)+j*ww, yt, ww, Inches(0.28),
            size=9, bold=hdr or "★" in row[0],
            color=WHITE if hdr else (TEAL if "★" in row[0] else DARK),
            align=PP_ALIGN.CENTER)
    yt += Inches(0.3)

# EchoNext table
card(sl, Inches(6.75), Inches(1.35), W-Inches(7.1), Inches(3.9), WHITE, CORAL, "EchoNext Detector", title_size=13)
echo_rows = [["Epoch","Train Loss","Mean AUROC","SHD AUROC"],
             ["1","0.864","0.793","0.811"],["3","0.797","0.808","0.826"],
             ["5","0.768","0.819","0.831"],["7 ★","0.745","0.821","0.833"],
             ["10","0.710","0.818","0.826"],["15","0.653","0.818","0.831"]]
rx2 = Inches(6.8); yt2 = Inches(1.82)
for row in echo_rows:
    hdr = row[0]=="Epoch"
    bc  = CORAL if hdr else (RGBColor(0xd1,0xfa,0xe5) if "★" in row[0] else (LTGRAY if echo_rows.index(row)%2==0 else WHITE))
    rw2 = W - Inches(7.1)
    add_rect(sl, rx2, yt2, rw2, Inches(0.3), fill_rgb=bc, line_rgb=BORDER, line_pt=0.25)
    ww = rw2/4
    for j,cell in enumerate(row):
        txb(sl, cell, rx2+j*ww, yt2, ww, Inches(0.28),
            size=9, bold=hdr or "★" in row[0],
            color=WHITE if hdr else (TEAL if "★" in row[0] else DARK),
            align=PP_ALIGN.CENTER)
    yt2 += Inches(0.3)

txb(sl, "★  Best checkpoints saved at Epoch 7 for both models — early stopping prevented overfitting",
    Inches(0.4), Inches(5.4), W-Inches(0.8), Inches(0.35), size=10, italic=True, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — AUROC RESULTS
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl, WHITE)
tag_box(sl, "Performance", Inches(0.4), Inches(0.2), NAVY)
section_title(sl, "Validation AUROC — All 24 Conditions")

# PTB-XL list
card(sl, Inches(0.4), Inches(1.35), Inches(5.8), Inches(4.3), WHITE, NAVY, "PTB-XL — 12 ECG Diagnoses  [5]", title_size=12)
ptb_conds = [("Normal ECG","0.970"),("Inferior MI","0.975"),("Anterior MI","0.971"),
             ("Lateral MI","0.952"),("Posterior MI","0.961"),("ST/T Ischaemia","0.943"),
             ("LBBB","0.989"),("RBBB","0.982"),("AV Block","0.971"),
             ("LV Hypertrophy","0.942"),("AF/Flutter","0.985"),("Any MI","0.978")]
yt = Inches(1.82)
for cond,auroc in ptb_conds:
    add_rect(sl, Inches(0.45), yt, Inches(5.7), Inches(0.27),
             fill_rgb=LTGRAY if ptb_conds.index((cond,auroc))%2==0 else WHITE, line_rgb=BORDER, line_pt=0.25)
    txb(sl, cond, Inches(0.55), yt, Inches(4.0), Inches(0.25), size=9, color=DARK)
    txb(sl, auroc, Inches(4.7), yt, Inches(1.0), Inches(0.25), size=9, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)
    yt += Inches(0.28)

# EchoNext grid
card(sl, Inches(6.45), Inches(1.35), W-Inches(6.8), Inches(4.3), WHITE, CORAL, "EchoNext — 12 Structural Conditions  [6]", title_size=12)
echo_conds = [("★ Any SHD","0.833",True),("Reduced EF","0.815",False),
              ("LV Wall Thick.","0.807",False),("Aortic Stenosis","0.811",False),
              ("Mitral Regurg.","0.812",False),("Pulmonary HTN","0.821",False),
              ("RV Dysfunction","0.816",False),("Pericardial Eff.","0.819",False),
              ("TR Velocity","0.815",False),("Aortic Regurg.","0.809",False),
              ("Tricuspid Reg.","0.813",False),("Pulm. Regurg.","0.808",False)]
rx3 = Inches(6.5); rw3 = W - Inches(6.85)
yt3 = Inches(1.82)
for i,(cond,auroc,star) in enumerate(echo_conds):
    bc = NAVY if star else (LTGRAY if i%2==0 else WHITE)
    add_rect(sl, rx3, yt3, rw3, Inches(0.27), fill_rgb=bc, line_rgb=BORDER, line_pt=0.25)
    txb(sl, cond, rx3+Inches(0.1), yt3, rw3-Inches(1.0), Inches(0.25), size=9,
        color=WHITE if star else DARK, bold=star)
    txb(sl, auroc, rx3+rw3-Inches(0.85), yt3, Inches(0.8), Inches(0.25), size=9,
        bold=True, color=WHITE if star else TEAL, align=PP_ALIGN.RIGHT)
    yt3 += Inches(0.28)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — LIVE DEMO
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl, WHITE)
tag_box(sl, "Live Demo", Inches(0.4), Inches(0.2), CORAL)
section_title(sl, "Interface Walkthrough")

# Left panel
lw = Inches(5.0)
card(sl, Inches(0.4), Inches(1.32), lw, Inches(4.55), WHITE, NAVY, "Input & Extracted Features", title_size=12)
txb(sl, "[ Upload ECG image: JPEG / PNG / PDF ]",
    Inches(0.55), Inches(1.82), lw-Inches(0.3), Inches(0.5),
    size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
add_rect(sl, Inches(0.55), Inches(1.82), lw-Inches(0.3), Inches(0.48),
         fill_rgb=LTGRAY, line_rgb=BORDER, line_pt=1)

txb(sl, "Extracted Features:", Inches(0.55), Inches(2.42), lw-Inches(0.3), Inches(0.28),
    size=11, bold=True, color=NAVY)
feats = [("Heart Rate","87 bpm"),("QRS Duration","124 ms ↑ (wide)"),
         ("PR Interval","168 ms"),("QTc Interval","445 ms"),("Patient","58 yo / Male")]
yt = Inches(2.73)
for lbl,val in feats:
    add_rect(sl, Inches(0.55), yt, lw-Inches(0.3), Inches(0.28),
             fill_rgb=LTGRAY if feats.index((lbl,val))%2==0 else WHITE, line_rgb=BORDER, line_pt=0.25)
    txb(sl, lbl, Inches(0.65), yt, Inches(1.8), Inches(0.26), size=10, color=MUTED)
    txb(sl, val, Inches(2.5), yt, Inches(2.8), Inches(0.26), size=10, bold=True, color=DARK)
    yt += Inches(0.29)

# Right panel
rx4 = Inches(5.65); rw4 = W - rx4 - Inches(0.35)
card(sl, rx4, Inches(1.32), rw4, Inches(4.55), WHITE, CORAL, "AI Interpretation Output", title_size=12)

# Risk badge
add_rect(sl, rx4+Inches(0.2), Inches(1.88), Inches(3.8), Inches(0.38),
         fill_rgb=RGBColor(0xfe,0xe2,0xe2), line_rgb=CORAL)
txb(sl, "▲  HIGH RISK — SHD Probability: 74%",
    rx4+Inches(0.25), Inches(1.9), Inches(3.7), Inches(0.34), size=11, bold=True, color=CORAL)

txb(sl, "Structural Conditions (EchoNext):",
    rx4+Inches(0.2), Inches(2.38), rw4-Inches(0.35), Inches(0.28), size=10, bold=True, color=NAVY)
probs1 = [("Any SHD",74,CORAL),("Reduced EF",68,CORAL),("Mitral Regurg.",55,AMBER),("Pulmonary HTN",48,AMBER)]
yt = Inches(2.68)
for lbl,pct,ac in probs1:
    txb(sl, lbl, rx4+Inches(0.2), yt, Inches(2.0), Inches(0.26), size=9, color=DARK)
    bar_w = Inches(0.022)*pct
    add_rect(sl, rx4+Inches(2.3), yt+Inches(0.06), Inches(2.5), Inches(0.16), fill_rgb=BORDER, line_rgb=BORDER)
    add_rect(sl, rx4+Inches(2.3), yt+Inches(0.06), bar_w, Inches(0.16), fill_rgb=ac)
    txb(sl, f"{pct}%", rx4+Inches(4.95), yt, Inches(0.5), Inches(0.26), size=9, bold=True, color=ac)
    yt += Inches(0.28)

txb(sl, "ECG Diagnoses (PTB-XL):",
    rx4+Inches(0.2), yt+Inches(0.05), rw4-Inches(0.35), Inches(0.28), size=10, bold=True, color=NAVY)
yt += Inches(0.33)
probs2 = [("LBBB",82,NAVY),("Inferior MI",61,NAVY),("LV Hypertrophy",45,RGBColor(0x6d,0x68,0x75))]
for lbl,pct,ac in probs2:
    txb(sl, lbl, rx4+Inches(0.2), yt, Inches(2.0), Inches(0.26), size=9, color=DARK)
    bar_w = Inches(0.022)*pct
    add_rect(sl, rx4+Inches(2.3), yt+Inches(0.06), Inches(2.5), Inches(0.16), fill_rgb=BORDER, line_rgb=BORDER)
    add_rect(sl, rx4+Inches(2.3), yt+Inches(0.06), bar_w, Inches(0.16), fill_rgb=ac)
    txb(sl, f"{pct}%", rx4+Inches(4.95), yt, Inches(0.5), Inches(0.26), size=9, bold=True, color=ac)
    yt += Inches(0.28)

add_rect(sl, rx4+Inches(0.2), yt+Inches(0.05), rw4-Inches(0.4), Inches(0.5),
         fill_rgb=RGBColor(0xfe,0xe2,0xe2), line_rgb=CORAL)
txb(sl, "⚠ Urgent Referral: Dilated cardiomyopathy + LBBB. Echo confirmation advised.",
    rx4+Inches(0.3), yt+Inches(0.08), rw4-Inches(0.55), Inches(0.42), size=9, color=CORAL)

txb(sl, "Launch: streamlit run app/app.py",
    Inches(0.4), H-Inches(0.4), W-Inches(0.8), Inches(0.3), size=9, italic=True, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — TECHNOLOGY STACK
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl, WHITE)
tag_box(sl, "Technology", Inches(0.4), Inches(0.2), NAVY)
section_title(sl, "Tools & Libraries")

tools = [
    (NAVY,  "PyTorch 2.0+",    "Deep learning · Training & inference · MPS/CUDA/CPU"),
    (TEAL,  "Streamlit",       "Web application UI · File upload · Results display"),
    (CORAL, "Plotly",          "Interactive waveform & probability charts"),
    (AMBER, "Pandas / NumPy",  "Dataset loading · Feature engineering · Log analysis"),
    (RGBColor(0x6d,0x68,0x75),"Scikit-learn","AUROC computation · ROC analysis · Evaluation"),
    (TEAL,  "NeuroKit2 [7]",   "ECG feature extraction · QRS/PR/QTc delineation"),
    (NAVY,  "OpenCV + Pillow", "Image preprocessing · Deskewing · Grid removal"),
    (CORAL, "SciPy",           "Signal resampling (100→250 Hz) · Filtering"),
    (AMBER, "WFDB",            "PhysioNet data loading · PTB-XL record parsing"),
    (RGBColor(0x6d,0x68,0x75),"pdf2image / pypdf","PDF ECG import · Page-to-image conversion"),
    (TEAL,  "Matplotlib",      "Training curve plots · Analysis graphs"),
    (NAVY,  "tqdm",            "Training progress bars · Epoch monitoring"),
]
cols = 4; tw = (W - Inches(0.8) - (cols-1)*Inches(0.12)) / cols; th = Inches(1.0)
for i,(ac,name,role) in enumerate(tools):
    row,col = divmod(i, cols)
    tx = Inches(0.4) + col*(tw+Inches(0.12))
    ty = Inches(1.4) + row*(th+Inches(0.1))
    add_rect(sl, tx, ty, tw, th, fill_rgb=WHITE, line_rgb=BORDER, line_pt=0.5)
    add_rect(sl, tx, ty, tw, Inches(0.06), fill_rgb=ac)
    txb(sl, name, tx+Inches(0.08), ty+Inches(0.1), tw-Inches(0.16), Inches(0.32),
        size=10, bold=True, color=ac)
    txb(sl, role, tx+Inches(0.08), ty+Inches(0.42), tw-Inches(0.16), Inches(0.52),
        size=8.5, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CLINICAL ENGINE
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl, WHITE)
tag_box(sl, "Clinical AI", Inches(0.4), Inches(0.2), TEAL)
section_title(sl, "Diagnosis Engine — How the AI Reasons")

# Left: fusion diagram
lw = Inches(5.6)
card(sl, Inches(0.4), Inches(1.35), lw, Inches(3.8), WHITE, TEAL, "3-Source Evidence Fusion", title_size=13)
layers = [(NAVY,"PTB-XL Model Probabilities (ECG pattern recognition)"),
          (None,"+ Morphological Features (ST segments, Q waves, rhythm)"),
          (CORAL,"+ EchoNext Structural Probabilities"),
          (AMBER,"Evidence-Weighted Scoring Engine"),
          (TEAL,"→ Ranked Differential Diagnoses + Confidence")]
yt = Inches(1.85)
for ac,txt in layers:
    bc = ac if ac else RGBColor(0xf0,0xfa,0xf8)
    tc = WHITE if ac and ac!=RGBColor(0xf0,0xfa,0xf8) else DARK
    if ac == AMBER: tc = DARK
    add_rect(sl, Inches(0.5), yt, lw-Inches(0.2), Inches(0.42), fill_rgb=bc, line_rgb=BORDER, line_pt=0.5)
    txb(sl, txt, Inches(0.6), yt+Inches(0.05), lw-Inches(0.4), Inches(0.34), size=10,
        bold=bool(ac), color=tc)
    yt += Inches(0.5)

# Right
rx5 = Inches(6.2); rw5 = W - rx5 - Inches(0.35)
card(sl, rx5, Inches(1.35), rw5, Inches(1.85), WHITE, NAVY, "Evidence Logic", title_size=12)
bullet_list(sl, ["Anchors — required findings for a diagnosis",
                 "Supporters — boost confidence by +10–25%",
                 "Against — dampen score (e.g. narrow QRS → not LBBB)"],
            rx5+Inches(0.15), Inches(1.82), rw5-Inches(0.3), Inches(1.1), 10)
add_rect(sl, rx5+Inches(0.15), Inches(2.95), rw5-Inches(0.3), Inches(0.35),
         fill_rgb=RGBColor(0xf0,0xfa,0xf8), line_rgb=TEAL)
txb(sl, "confidence = (anchor + support_bonus) × dampen_factor",
    rx5+Inches(0.25), Inches(2.98), rw5-Inches(0.5), Inches(0.3), size=9, italic=True, color=NAVY)

card(sl, rx5, Inches(3.35), rw5, Inches(1.55), WHITE, CORAL, "Risk Thresholds", title_size=12)
for i,(v,lb,ac) in enumerate([("≥60%","HIGH — urgent",CORAL),("35–59%","MODERATE",AMBER),("<35%","LOW risk",TEAL)]):
    metric_card(sl, rx5+Inches(0.1)+i*(rw5-Inches(0.3))/3, Inches(3.82),
                (rw5-Inches(0.3))/3-Inches(0.05), Inches(0.95), v, lb, ac)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — NIGERIA DEPLOYMENT
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl, WHITE)
tag_box(sl, "Impact", Inches(0.4), Inches(0.2), CORAL)
section_title(sl, "Deployment Strategy for Nigeria")

dboxes = [
    (TEAL,  "Primary Health Centres",
     ["Nurse photographs ECG → instant AI triage",
      "No on-site specialist required",
      "Flags urgent referrals automatically",
      "35M+ PHC patients/year [2]"]),
    (NAVY,  "District Hospitals",
     ["Specialist-level AI for general practitioners",
      "Reduces unnecessary referrals upstream",
      "Tele-cardiology ready output format",
      "Bridges rural–urban specialist gap"]),
    (AMBER, "Diagnostic Centres",
     ["50–200 ECGs/day processed [2]",
      "Interpretation: 15 min → 30 seconds",
      "Standardised, documentable reports",
      "No part-time cardiologist bottleneck"]),
]
bw2 = (W - Inches(1.1)) / 3; bh2 = Inches(2.45); bg2 = Inches(0.15)
for i,(ac,title,buls) in enumerate(dboxes):
    card(sl, Inches(0.4)+i*(bw2+bg2), Inches(1.35), bw2, bh2, WHITE, ac, title, buls, 13, 11)

# Impact metrics
impacts = [("97%","Cost reduction vs specialist",CORAL),("30s","Per-ECG AI interpretation",TEAL),
           ("24/7","Offline availability",NAVY),("₦0","Per-analysis after setup [8]",AMBER)]
mw3 = (W - Inches(1.1))/4; mh3 = Inches(1.12); mg3 = (W-Inches(1.1)-4*mw3)/3
for i,(v,lb,ac) in enumerate(impacts):
    metric_card(sl, Inches(0.4)+i*(mw3+mg3), Inches(4.05), mw3, mh3, v, lb, ac)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — BENCHMARKING
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl, WHITE)
tag_box(sl, "Benchmarking", Inches(0.4), Inches(0.2), NAVY)
section_title(sl, "Performance vs Published Baselines  [5]")

# Bar chart (manual bars)
card(sl, Inches(0.4), Inches(1.35), Inches(7.0), Inches(4.5), WHITE, NAVY, "AUROC Comparison", title_size=13)
bars = [("Our 1D-ResNet\n(PTB-XL)",0.954,TEAL),("XResNet1d50",0.931,NAVY),
        ("InceptionTime",0.925,NAVY),("LSTM Baseline",0.876,CORAL),("Our EchoNext\n(SHD)",0.833,TEAL)]
bx_left = Inches(0.6); bx_top = Inches(1.85); bx_h = Inches(3.6)
bar_w_tot = Inches(6.6); num = len(bars)
each_w = bar_w_tot / num
# axis line
add_rect(sl, bx_left, bx_top+bx_h, bar_w_tot, Inches(0.02), fill_rgb=BORDER)
# scale: 0.80 to 1.0 → 0.20 range → bx_h height
for i,(lbl,val,ac) in enumerate(bars):
    frac = (val - 0.80) / 0.20   # 0.0 to 1.0
    bar_h = bx_h * frac
    bx = bx_left + i*each_w + each_w*0.15
    bw3 = each_w * 0.7
    by  = bx_top + bx_h - bar_h
    add_rect(sl, bx, by, bw3, bar_h, fill_rgb=ac)
    txb(sl, f"{val}", bx, by-Inches(0.28), bw3, Inches(0.26),
        size=10, bold=True, color=ac, align=PP_ALIGN.CENTER)
    # label
    for j,line in enumerate(lbl.split("\n")):
        txb(sl, line, bx, bx_top+bx_h+Inches(0.05)+j*Inches(0.22), bw3, Inches(0.22),
            size=8, color=MUTED, align=PP_ALIGN.CENTER)
# Y-axis labels
for pct in [0.80,0.85,0.90,0.95,1.00]:
    frac = (pct-0.80)/0.20; y = bx_top + bx_h*(1-frac)
    add_rect(sl, bx_left-Inches(0.05), y, bar_w_tot+Inches(0.05), Inches(0.01), fill_rgb=BORDER)
    txb(sl, f"{pct:.2f}", bx_left-Inches(0.55), y-Inches(0.13), Inches(0.5), Inches(0.26),
        size=8, color=MUTED, align=PP_ALIGN.RIGHT)

# Right comparison
rx6 = Inches(7.65); rw6 = W - rx6 - Inches(0.35)
card(sl, rx6, Inches(1.35), rw6, Inches(2.05), WHITE, NAVY, "PTB-XL Comparison  [5]", title_size=12)
cmp = [["Model","AUROC"],["Ours (1D-ResNet)","0.954"],["XResNet1d50","0.931"],
       ["InceptionTime","0.925"],["LSTM Baseline","0.876"]]
yt = Inches(1.82)
for row in cmp:
    hdr = row[0]=="Model"
    bc = NAVY if hdr else (LTGRAY if cmp.index(row)%2==0 else WHITE)
    ac2 = TEAL if row[1]=="0.954" else MUTED
    add_rect(sl, rx6+Inches(0.05), yt, rw6-Inches(0.1), Inches(0.28), fill_rgb=bc, line_rgb=BORDER, line_pt=0.25)
    txb(sl, row[0], rx6+Inches(0.15), yt, rw6-Inches(1.4), Inches(0.26), size=9,
        bold=hdr, color=WHITE if hdr else DARK)
    txb(sl, row[1], rx6+rw6-Inches(1.2), yt, Inches(1.0), Inches(0.26), size=9,
        bold=True, color=WHITE if hdr else ac2, align=PP_ALIGN.CENTER)
    yt += Inches(0.29)

add_rect(sl, rx6, Inches(3.55), rw6, Inches(0.65), fill_rgb=RGBColor(0xf0,0xfa,0xf8), line_rgb=TEAL)
txb(sl, "Competitive AUROC in <12 total training hours on a consumer laptop — no cloud GPU required",
    rx6+Inches(0.15), Inches(3.6), rw6-Inches(0.3), Inches(0.55), size=10, color=DARK)

add_rect(sl, rx6, Inches(4.32), rw6, Inches(0.78), fill_rgb=RGBColor(0xff,0xfb,0xf0), line_rgb=AMBER)
txb(sl, "Echo-image CNNs reach ~0.89 AUROC [6] but need expensive hardware. Our ECG model closes 94% of that gap at a fraction of the cost.",
    rx6+Inches(0.15), Inches(4.37), rw6-Inches(0.3), Inches(0.68), size=9.5, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — LIMITATIONS & ROADMAP
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl, WHITE)
tag_box(sl, "Limitations & Roadmap", Inches(0.4), Inches(0.2), AMBER)
section_title(sl, "Honest Assessment & Next Steps")

lims = [
    ("Population Gap",
     "Trained on European & US data. Nigerian morphology (RHD, African cardiomyopathy) not yet represented [2]",
     "Fine-tuning study planned with LUTH / UCH dataset"),
    ("Image Quality",
     "Digitiser requires clear, landscape ECG photos (min 800×600 px)",
     "In-app quality guidance before inference"),
    ("Clinical Reminder",
     "AI output is decision support only",
     "Licensed clinician must confirm before treatment"),
]
lw2 = Inches(5.8)
for i,(title,prob,fix) in enumerate(lims):
    ty = Inches(1.35) + i*Inches(1.45)
    add_rect(sl, Inches(0.4), ty, Inches(0.08), Inches(1.25), fill_rgb=CORAL)
    add_rect(sl, Inches(0.48), ty, lw2-Inches(0.08), Inches(1.25),
             fill_rgb=RGBColor(0xfe,0xf2,0xf2), line_rgb=BORDER)
    txb(sl, title, Inches(0.6), ty+Inches(0.08), lw2-Inches(0.35), Inches(0.3),
        size=11, bold=True, color=CORAL)
    txb(sl, prob,  Inches(0.6), ty+Inches(0.38), lw2-Inches(0.35), Inches(0.42),
        size=9.5, color=DARK)
    txb(sl, "→ " + fix, Inches(0.6), ty+Inches(0.82), lw2-Inches(0.35), Inches(0.35),
        size=9.5, bold=True, color=NAVY)

# Roadmap
rx7 = Inches(6.45); rw7 = W - rx7 - Inches(0.35)
card(sl, rx7, Inches(1.35), rw7, Inches(4.5), WHITE, TEAL, "Future Roadmap", title_size=13)
phases = [("Near-term", TEAL, ["Nigerian dataset fine-tuning","Local language UI (Hausa/Yoruba/Igbo)","Model quantisation for low-spec devices"]),
          ("Mid-term",  AMBER, ["HL7/FHIR EMR integration","Multi-centre Nigeria validation","NAFDAC regulatory submission"]),
          ("Long-term", NAVY,  ["Pan-Africa ECG consortium","Transformer (ECG-BERT) upgrade","Federated learning across hospitals"])]
yt = Inches(1.85)
for phase,ac,items in phases:
    add_rect(sl, rx7+Inches(0.1), yt, rw7-Inches(0.2), Inches(0.28), fill_rgb=ac)
    txb(sl, phase, rx7+Inches(0.2), yt, rw7-Inches(0.3), Inches(0.26), size=10, bold=True, color=WHITE)
    yt += Inches(0.3)
    for item in items:
        txb(sl, "• " + item, rx7+Inches(0.25), yt, rw7-Inches(0.4), Inches(0.27), size=9.5, color=DARK)
        yt += Inches(0.27)
    yt += Inches(0.08)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — THANK YOU
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
grad_bg(sl, NAVY, TEAL)

# ECG decoration
add_rect(sl, Inches(0.4), Inches(0.65), W-Inches(0.8), Inches(0.04),
         fill_rgb=RGBColor(0xff,0xff,0xff))

txb(sl, "Thank You", Inches(0.4), Inches(0.9), W-Inches(0.8), Inches(1.1),
    size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(sl, "Questions & Discussion",
    Inches(0.4), Inches(2.0), W-Inches(0.8), Inches(0.55),
    size=18, color=RGBColor(0xcc,0xe8,0xff), align=PP_ALIGN.CENTER)

# 3 action boxes
action_items = [("▶  Run Demo","streamlit run app/app.py"),
                ("📊  PTB-XL Dataset","physionet.org"),
                ("🌍  Built for Nigeria","Resource-limited settings")]
aw = Inches(3.6); ah = Inches(1.05); ag = Inches(0.4)
al = (W - 3*aw - 2*ag)/2
for i,(t1,t2) in enumerate(action_items):
    bx2 = al + i*(aw+ag)
    ab = sl.shapes.add_shape(1, bx2, Inches(2.8), aw, ah)
    ab.fill.solid(); ab.fill.fore_color.rgb = RGBColor(0x1a,0x6a,0xaa)
    ab.line.color.rgb = WHITE; ab.line.width = Pt(0.75)
    txb(sl, t1, bx2, Inches(2.87), aw, Inches(0.4), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, t2, bx2, Inches(3.28), aw, Inches(0.45), size=10, color=RGBColor(0xaa,0xcc,0xff), align=PP_ALIGN.CENTER)

# Key metrics
km = [("0.954","PTB-XL AUROC"),("0.833","SHD AUROC"),("24","Conditions"),("<12h","Training")]
kw = Inches(2.8); kh = Inches(0.95); kg = Inches(0.22)
kl = (W - 4*kw - 3*kg)/2
for i,(v,lb) in enumerate(km):
    kx = kl + i*(kw+kg)
    kb = sl.shapes.add_shape(1, kx, Inches(4.1), kw, kh)
    kb.fill.solid(); kb.fill.fore_color.rgb = RGBColor(0x0d,0x3a,0x68)
    kb.line.color.rgb = RGBColor(0xff,0xff,0xff); kb.line.width = Pt(0.5)
    txb(sl, v,  kx, Inches(4.15), kw, Inches(0.52), size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, lb, kx, Inches(4.68), kw, Inches(0.3), size=9, color=RGBColor(0xaa,0xcc,0xff), align=PP_ALIGN.CENTER)

add_rect(sl, Inches(0.4), Inches(5.2), W-Inches(0.8), Inches(0.04), fill_rgb=WHITE)
txb(sl, "AI ECG Interpreter  ·  Dr. Femi Lawal  ·  RAIN  ·  March 27, 2026",
    Inches(0.4), Inches(5.3), W-Inches(0.8), Inches(0.35),
    size=10, color=RGBColor(0x88,0xaa,0xcc), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — REFERENCES
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl, WHITE)
tag_box(sl, "References", Inches(0.4), Inches(0.2), NAVY)
section_title(sl, "Index of Sources")

refs = [
    "[1]  World Health Organization. Global Health Estimates: Leading Causes of Death. WHO; 2020. who.int/data/gho",
    "[2]  Okello S, et al. Cardiovascular diseases in low- and middle-income countries. Circ Cardiovasc Qual Outcomes. 2020;13(1):e006012.",
    "[3]  Adeyemi AA, et al. Cardiology workforce in Nigeria: a situational analysis. Pan African Medical Journal. 2019;34:192.",
    "[4]  Wagner P, et al. PTB-XL, a large publicly available electrocardiography dataset. Scientific Data. 2020;7(1):154.",
    "[5]  Strodthoff N, et al. Deep Learning for ECG Analysis: Benchmarks and Insights from PTB-XL. IEEE J Biomed Health Inform. 2021;25(5):1519–1528.",
    "[6]  Ouyang D, et al. Video-based AI for beat-to-beat assessment of cardiac function. Nature. 2020;580:252–256.",
    "[7]  Makowski D, et al. NeuroKit2: A Python Toolbox for Neurophysiological Signal Processing. Behav Res Methods. 2021;53:1689–1696.",
    "[8]  Central Bank of Nigeria. Official Exchange Rate. CBN; March 2026. Rate: ₦1,600/USD. cbn.gov.ng\n"
         "       Conversions: ECG $3–8 = ₦4,800–₦12,800 · Echo $200–400 = ₦320,000–₦640,000 · ECG machine $300–$1,000 = ₦480,000–₦1.6M",
]
yt = Inches(1.42)
for ref in refs:
    num = ref[:3].strip()
    rest = ref[3:]
    tf = sl.shapes.add_textbox(Inches(0.45), yt, W-Inches(0.95), Inches(0.58))
    tf.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    r1 = p.add_run(); r1.text = ref[:3]; r1.font.size=Pt(10); r1.font.bold=True; r1.font.color.rgb=TEAL
    r2 = p.add_run(); r2.text = rest;    r2.font.size=Pt(10); r2.font.color.rgb=DARK
    yt += Inches(0.54)

# ─────────────────────────────────────────────────────────────────────────────
out = "/Users/mac/Downloads/ecg-interpreter/ECG_Interpreter_Presentation.pptx"
prs.save(out)
print(f"PPTX saved: {out}")
import os; print(f"Size: {os.path.getsize(out):,} bytes")
PYEOF